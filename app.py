import streamlit as st
import sqlite3
import smtplib
import ssl
import uuid
import re
from email.message import EmailMessage
from datetime import datetime
from hashlib import sha256
import os
import json
import requests
import tempfile
import urllib.parse
import PyPDF2
import docx
import pyttsx3
import speech_recognition as sr
from bs4 import BeautifulSoup
import random
from PIL import Image
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from fpdf import FPDF
import base64
import subprocess
import feedparser
import sys
import io
import traceback
from contextlib import redirect_stdout, redirect_stderr
import threading
import sounddevice as sd
from scipy.io.wavfile import write, read
import numpy as np
from gtts import gTTS
from deep_translator import GoogleTranslator
import matplotlib.pyplot as plt
import requests as req
from openai import OpenAI
from groq import Groq
import graphviz
import pytesseract
from bs4 import Tag
import toml

# Load secrets from secrets.toml
#secrets = toml.load("secrets.toml")

# === API KEYS ===
OPENROUTER_API_KEY = st.secrets["openrouter"]["token"]
GITHUB_API_TOKEN   = st.secrets["github"]["token"]
GROQ_API_KEY       = st.secrets["groq"]["token"]
A4F_API_KEY        = st.secrets["a4f"]["token"]
WEATHER_API_KEY    = st.secrets["weather"]["token"]
EMAIL_SENDER       = st.secrets["email"]["sender"]
EMAIL_PASSWORD     = st.secrets["email"]["password"]
WEBSEARCH_API_KEY  = st.secrets["websearch"]["token"]
# === Model Configurations ===
model_sources = {
    "GitHub GPT-4.1": {
        "type": "github",
        "model": "openai/gpt-4.1",
        "base_url": "https://models.github.ai/inference",
        "api_key": secrets["github"]["token"]
    },
    "Groq Llama 4 Maverick": {
        "type": "groq",
        "model": "meta-llama/llama-4-maverick-17b-128e-instruct",
        "api_key": secrets["groq"]["token"]
    },
    "OpenRouter DeepSeek Chat v3": {
        "type": "openrouter",
        "model": "deepseek/deepseek-chat-v3-0324:free",
        "base_url": "https://openrouter.ai/api/v1",
        "api_key": secrets["openrouter"]["token"]
    },
    "OpenRouter DeepSeek R1": {
        "type": "openrouter",
        "model": "deepseek/deepseek-r1:free",
        "base_url": "https://openrouter.ai/api/v1",
        "api_key": secrets["openrouter"]["token"]
    },
    "A4F Qwen3-235B": {
        "type": "a4f",
        "model": "provider-5/Qwen/Qwen3-235B-A22B",
        "base_url": "https://api.a4f.co/v1",
        "api_key": secrets["a4f"]["token"]
    },
    "A4F Grok-4-0709": {
        "type": "a4f",
        "model": "provider-3/grok-4-0709",
        "base_url": "https://api.a4f.co/v1",
        "api_key": secrets["a4f"]["token"]
    },
    "A4F Gemini-2.5-Flash": {
        "type": "a4f",
        "model": "provider-5/gemini-2.5-flash-preview-04-17",
        "base_url": "https://api.a4f.co/v1",
        "api_key": secrets["a4f"]["token"]
    }
}

# === AI Fallback Function ===
def get_ai_response(messages, model_preference=None):
    """
    Try making API call to preferred model, fallback if fails, return response and model used.
    """
    model_try_order = []
    if model_preference and model_preference in model_sources:
        model_try_order.append(model_preference)
    for m in model_sources:
        if m not in model_try_order:
            model_try_order.append(m)

    for model_name in model_try_order:
        selected_model = model_sources.get(model_name)
        if selected_model is None:
            continue
        try:
            if selected_model["type"] == "openrouter":
                headers = {
                    "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                    "Content-Type": "application/json",
                    "HTTP-Referer": "https://yourdomain.com", 
                    "X-Title": "OpenRouter Fallback"
                }
                payload = {
                    "model": selected_model["model"],
                    "messages": messages
                }
                response = req.post("https://openrouter.ai/api/v1/chat/completions", 
                                    headers=headers, json=payload, timeout=30)
                if response.status_code == 200:
                    data = response.json()
                    return data["choices"][0]["message"]["content"], model_name
                else:
                    continue

            elif selected_model["type"] == "github":
                github_client = OpenAI(
                    base_url="https://models.github.ai/inference", 
                    api_key=GITHUB_API_TOKEN
                )
                response = github_client.chat.completions.create(
                    model=selected_model["model"],
                    messages=messages,
                    temperature=1.0,
                    top_p=1.0
                )
                if response and response.choices and len(response.choices) > 0:
                    return response.choices[0].message.content, model_name
                else:
                    continue

            elif selected_model["type"] == "groq":
                groq_client = Groq(api_key=GROQ_API_KEY)
                response = groq_client.chat.completions.create(
                    model=selected_model["model"],
                    messages=messages,
                    temperature=1.0,
                    max_completion_tokens=1024,
                    top_p=1.0
                )
                if response and response.choices and len(response.choices) > 0:
                    return response.choices[0].message.content, model_name
                else:
                    continue

            elif selected_model["type"] == "a4f":
                a4f_client = OpenAI(
                    api_key=A4F_API_KEY,
                    base_url="https://api.a4f.co/v1"
                )
                response = a4f_client.chat.completions.create(
                    model=selected_model["model"],
                    messages=messages
                )
                if response and response.choices and len(response.choices) > 0:
                    return response.choices[0].message.content, model_name
                else:
                    continue

        except Exception as e:
            continue

    return "[❌ All model APIs failed or rate-limited.]", "None"

# --------- DB Setup ---------
conn = sqlite3.connect("chatbot_local.db", check_same_thread=False)
c = conn.cursor()
c.execute('''CREATE TABLE IF NOT EXISTS users (
    email TEXT PRIMARY KEY,
    password TEXT,
    is_verified INTEGER,
    verification_code TEXT
)''')
c.execute('''CREATE TABLE IF NOT EXISTS chat_history (
    id TEXT PRIMARY KEY,
    email TEXT,
    user_msg TEXT,
    bot_msg TEXT,
    timestamp TEXT
)''')
conn.commit()

# --------- SMTP Email Config ---------
SMTP_SERVER = 'smtp.gmail.com'
SMTP_PORT = 465

# --------- Helper Functions ---------
def send_verification_email(to_email, code):
    msg = EmailMessage()
    msg['Subject'] = 'Your OTP Verification Code'
    msg['From'] = EMAIL_SENDER
    msg['To'] = to_email
    msg.set_content(f"Your verification code is: {code}")
    context = ssl.create_default_context()
    with smtplib.SMTP_SSL(SMTP_SERVER, SMTP_PORT, context=context) as server:
        server.login(EMAIL_SENDER, EMAIL_PASSWORD)
        server.send_message(msg)

def validate_email(email):
    return re.match(r"[^@]+@[^@]+\.[^@]+", email)

def hash_password(password):
    return sha256(password.encode()).hexdigest()

def is_strong_password(password):
    return (len(password) >= 8 and
            re.search(r'[A-Z]', password) and
            re.search(r'[a-z]', password) and
            re.search(r'[0-9]', password) and
            re.search(r'[^A-Za-z0-9]', password))

# -------------- Setup --------------
SAVE_DIR = "saved_chats"
os.makedirs(SAVE_DIR, exist_ok=True)
st.set_page_config(page_title="🧠 AI Chatbot All-in-One", layout="wide")

# ====== 🪐 Space & Modern UI Effects Injection ======
st.markdown('''
<style>
.gradient-btn, button[data-testid="baseButton"], .stButton>button {
    background: rgba(255,255,255,0.13);
    background: linear-gradient(270deg, rgba(162,89,255,0.33), rgba(0,212,255,0.22), rgba(255,110,196,0.18), rgba(162,89,255,0.33));
    background-size: 600% 600%;
    color: #fff !important;
    border: none !important;
    border-radius: 16px !important;
    font-weight: 700;
    box-shadow: 0 2px 12px #a259ff44, 0 4px 24px #00d4ff22;
    transition: box-shadow 0.3s, transform 0.2s;
    animation: gradientMove 8s ease-in-out infinite;
    backdrop-filter: blur(8px) saturate(1.2);
    outline: none !important;
}
.gradient-btn:hover, button[data-testid="baseButton"]:hover, .stButton>button:hover {
    box-shadow: 0 6px 24px #a259ff88, 0 4px 24px #00d4ff55;
    transform: translateY(-2px) scale(1.04);
    background: linear-gradient(90deg,#5f5fff,#b68cff,#00d4ff,#ff6ec4,#a259ff); 
    background-size: 600% 600%;
}
.glass-card {
    background: rgba(30,16,64,0.35) !important;
    border: 2px solid #a259ff44 !important;
    border-radius: 24px !important;
    box-shadow: 0 8px 32px 0 #a259ff22, 0 2px 8px 0 #0002;
    backdrop-filter: blur(8px) saturate(1.1);
    transition: box-shadow 0.3s, border 0.3s;
    padding: 1.1rem 1.6rem 1.1rem 1.6rem;
    margin-bottom: 1.2rem;
}
.glass-card:hover {
    border: 2.5px solid #a259ffbb !important;
    box-shadow: 0 12px 36px 0 #a259ff55, 0 2px 8px 0 #0003;
}
.info-icon {
    display: inline-block;
    margin-left: 6px;
    color: #a259ff;
    background: rgba(162,89,255,0.14);
    border-radius: 50%;
    width: 18px;
    height: 18px;
    text-align: center;
    font-size: 14px;
    cursor: pointer;
    position: relative;
}
.info-icon:hover .tooltip {
    display: block;
}
.tooltip {
    display: none;
    position: absolute;
    left: 22px;
    top: -2px;
    z-index: 99;
    background: rgba(30,16,64,0.92);
    color: #fff;
    padding: 8px 14px;
    border-radius: 10px;
    font-size: 0.98rem;
    box-shadow: 0 2px 12px #a259ff44;
    white-space: pre-line;
    min-width: 120px;
    max-width: 320px;
}
@keyframes gradientMove {
    0% {background-position:0% 50%}
    50% {background-position:100% 50%}
    100% {background-position:0% 50%}
}
</style>
''', unsafe_allow_html=True)

def info_icon(tooltip_text, key=None):
    """Render an info icon with a tooltip and return HTML as string."""
    icon_html = f'''<span class="info-icon" tabindex="0">&#9432;<span class="tooltip">{tooltip_text}</span></span>'''
    return icon_html

st.markdown('''
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;700&family=Montserrat:wght@700&family=Poppins:wght@600&display=swap');
html, body, [class*="css"] {
    font-family: 'Poppins', 'Montserrat', 'Inter', sans-serif !important;
    background: #0e1021 !important;
    color: #ececff;
    min-height: 100vh;
    overflow-x: hidden;
}
/* Nebula background */
body::before {
    content: "";
    position: fixed;
    z-index: -3;
    top: 0; left: 0; width: 100vw; height: 100vh;
    background: radial-gradient(ellipse at 60% 40%, #6e00ff88 0%, #00d4ff44 50%, #0e1021 100%);
    animation: nebulaMove 24s linear infinite alternate;
    filter: blur(8px) brightness(1.1) saturate(1.3);
}
@keyframes nebulaMove {
    0% { transform: scale(1) rotate(0deg); }
    100% { transform: scale(1.2) rotate(8deg); }
}
/* Glassmorphism card */
.stContainer, .glass-card {
    background: rgba(30,16,64,0.35) !important;
    border: 2px solid #a259ff44 !important;
    border-radius: 24px !important;
    box-shadow: 0 8px 32px 0 #a259ff22, 0 2px 8px 0 #0002;
    backdrop-filter: blur(8px) saturate(1.1);
    transition: box-shadow 0.3s, border 0.3s;
}
.stContainer:hover, .glass-card:hover {
    border: 2.5px solid #a259ffbb !important;
    box-shadow: 0 12px 36px 0 #a259ff55, 0 2px 8px 0 #0003;
}
/* Gradient animated button */
.gradient-btn, button[data-testid="baseButton"] {
    background: linear-gradient(270deg, #a259ff, #00d4ff, #ff6ec4, #a259ff);
    background-size: 600% 600%;
    color: #fff !important;
    border: none !important;
    border-radius: 16px !important;
    font-weight: 700;
    box-shadow: 0 2px 12px #a259ff44;
    transition: box-shadow 0.3s, transform 0.2s;
    animation: gradientMove 8s ease-in-out infinite;
}
.gradient-btn:hover, button[data-testid="baseButton"]:hover {
    box-shadow: 0 6px 24px #a259ff88;
    transform: translateY(-2px) scale(1.04);
}
@keyframes gradientMove {
    0% {background-position:0% 50%}
    50% {background-position:100% 50%}
    100% {background-position:0% 50%}
}
/* Pulse glow for interactive elements */
.pulse-glow, .gradient-btn:active, button[data-testid="baseButton"]:active {
    box-shadow: 0 0 16px 4px #a259ffcc, 0 2px 8px #0002;
    animation: pulse 1.2s infinite alternate;
}
@keyframes pulse {
    from { box-shadow: 0 0 8px 0 #a259ff88; }
    to { box-shadow: 0 0 24px 8px #00d4ff88; }
}
/* Gradient text */
.gradient-text, .stTitle, .stMarkdown h1, h1, h2, h3 {
    background: linear-gradient(90deg, #a259ff, #00d4ff, #ff6ec4, #a259ff);
    background-size: 300% 300%;
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    animation: gradientMove 8s ease-in-out infinite;
    font-family: 'Montserrat', 'Poppins', 'Inter', sans-serif !important;
}
/* Smooth transitions */
*, *:before, *:after {
    transition: all 0.25s cubic-bezier(.4,2,.6,1) !important;
}
</style>
<!-- Starfield and floating particles canvas -->
<canvas id="starfield-bg" style="position:fixed;top:0;left:0;width:100vw;height:100vh;z-index:-2;pointer-events:none;"></canvas>
<script>
const canvas = document.getElementById('starfield-bg');
const ctx = canvas.getContext('2d');
let w = window.innerWidth, h = window.innerHeight;
canvas.width = w; canvas.height = h;
let stars = Array.from({length:120},()=>({
    x:Math.random()*w, y:Math.random()*h, z:Math.random()*w, o:Math.random(),
    r:Math.random()*1.1+0.3, speed:Math.random()*0.7+0.2, angle:Math.random()*6.28
}));
let particles = Array.from({length:16},()=>({
    x:Math.random()*w, y:Math.random()*h, vx:Math.random()*2-1, vy:Math.random()*2-1, r:Math.random()*8+7, a:Math.random()*360
}));
function drawStarfield(){
    ctx.clearRect(0,0,w,h);
    // Nebula overlay
    let grad = ctx.createRadialGradient(w/2,h/2,Math.min(w,h)/8,w/2,h/2,Math.max(w,h)/1.1);
    grad.addColorStop(0,"#6e00ff33"); grad.addColorStop(0.3,"#00d4ff22"); grad.addColorStop(1,"#0e1021");
    ctx.fillStyle = grad; ctx.fillRect(0,0,w,h);
    // Rotating stars
    let time = Date.now()/12000;
    for(let s of stars){
        let angle = s.angle + time;
        let x = w/2 + (s.x-w/2)*Math.cos(angle)-(s.y-h/2)*Math.sin(angle);
        let y = h/2 + (s.x-w/2)*Math.sin(angle)+(s.y-h/2)*Math.cos(angle);
        ctx.globalAlpha = s.o*0.7+0.3;
        ctx.beginPath(); ctx.arc(x,y,s.r,0,2*Math.PI); ctx.fillStyle="#fff"; ctx.fill();
    }
    // Floating particles
    for(let p of particles){
        ctx.save();
        ctx.translate(p.x,p.y);
        ctx.rotate((p.a+time*50)*Math.PI/180);
        ctx.globalAlpha = 0.7;
        ctx.beginPath();
        ctx.arc(0,0,p.r,0,2*Math.PI);
        ctx.fillStyle = ctx.createRadialGradient(0,0,0,0,0,p.r);
        ctx.fillStyle.addColorStop(0,"#fff");
        ctx.fillStyle.addColorStop(0.5,"#a259ff");
        ctx.fillStyle.addColorStop(1,"#00d4ff00");
        ctx.fill();
        ctx.restore();
        p.x += p.vx; p.y += p.vy; p.vx *= 0.99; p.vy *= 0.99; p.a += 0.2;
        if(p.x<0||p.x>w) p.vx*=-1;
        if(p.y<0||p.y>h) p.vy*=-1;
    }
    ctx.globalAlpha = 1.0;
    requestAnimationFrame(drawStarfield);
}
drawStarfield();
window.addEventListener('resize',()=>{w=window.innerWidth;h=window.innerHeight;canvas.width=w;canvas.height=h;});
</script>
''', unsafe_allow_html=True)
# ====== END Space & Modern UI Effects Injection ======
st.markdown('''
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;700&display=swap');
    html, body, [class*="css"]  {
        font-family: 'Inter', sans-serif !important;
        background: #181818 !important;
        color: #232323;
    }
    .chat-container {
        background: #232323ee;
        border-radius: 20px;
        box-shadow: 0 4px 32px 0 #00000044;
        padding: 20px 14px 20px 14px;
        margin-bottom: 10px;
        max-height: 520px;
        overflow-y: auto;
        position: relative;
    }
    .msg-left {
        background: #fff;
        color: #232323;
        border-radius: 16px 16px 16px 6px;
        padding: 13px 18px;
        margin: 10px 0;
        max-width: 70%;
        float: left;
        clear: both;
        box-shadow: 0 2px 12px 0 #00000020;
        font-size: 1.08rem;

    }
    .msg-right {
        background: #fff;
        color: #232323;
        border-radius: 16px 16px 6px 16px;
        padding: 13px 18px;
        margin: 10px 0;
        max-width: 70%;
        float: right;
        clear: both;
        text-align: right;
        box-shadow: 0 2px 12px 0 #00000020;
        font-size: 1.08rem;

    }
    .stChatInputContainer, .chat-input {
        background: #fff !important;
        border-radius: 20px !important;
        box-shadow: 0 2px 16px 0 #00000030;
        padding: 10px 18px !important;
        border: 1.5px solid #e6e6e6 !important;
        display: flex;
        align-items: center;
        margin-top: 10px;
        margin-bottom: 0;
        width: 100% !important;
        position: static !important;
        left: unset !important;
        bottom: unset !important;
        transform: none !important;
    }
    .stChatInputContainer input, .chat-input input {
        background: transparent;

        color: #232323;
        font-size: 1.1rem;
        flex: 1;
        outline: none;
    }
    .stChatInputContainer button, .chat-input button {
        background: #232323;
        color: #fff;

        border-radius: 50%;
        box-shadow: 0 0 8px 1px #00000033;
        width: 40px;
        height: 40px;
        margin-left: 10px;
        font-size: 1.3rem;
        cursor: pointer;
        transition: box-shadow 0.2s, background 0.2s;
    }
    .stChatInputContainer button:hover, .chat-input button:hover {
        box-shadow: 0 0 16px 4px #23232366;
        background: #444;
    }
    .stButton > button {
        background: #fff;
        color: #232323;
        border-radius: 14px;
        box-shadow: 0 0 6px 1px #00000022;

        font-weight: 600;
        font-size: 1.04rem;
        transition: box-shadow 0.2s, background 0.2s;
    }
    .stButton > button:hover {
        box-shadow: 0 0 12px 3px #23232333;
        background: #f2f2f2;
    }
    .stTabs [data-baseweb="tab"] {
        background: #232323;
        color: #fff;
        border-radius: 10px 10px 0 0;
        font-weight: 600;
        font-size: 1.04rem;
        margin-right: 4px;

        box-shadow: 0 2px 8px 0 #00000010;
    }
    .stTabs [aria-selected="true"] {
        background: #fff;
        color: #232323;
    }
    .scroll-bottom { height: 1px; }
    ::selection { background: #23232322; }
    </style>
''', unsafe_allow_html=True)

# === Space Earth Horizon, Moving Particles, Spotlight, and Fading Text ===
st.markdown(
    '''
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@400;600;900&family=Montserrat:wght@400;700&display=swap');
    html, body, .stApp {
        height: 100%;
        margin: 0;
        font-family: 'Poppins', 'Montserrat', 'Inter', sans-serif;
        background: #101010;
        overflow-x: hidden;
    }
    /* Main space background video */
    #bgvid {
        position: fixed;
        top: 0; left: 0;
        min-width: 100vw;
        min-height: 100vh;
        width: 100vw; height: 100vh;
        object-fit: cover;
        z-index: 0;
        opacity: 0.18;
        pointer-events: none;
        filter: blur(0.5px) brightness(1.1) saturate(1.1);
    }
    /* Earth horizon video at bottom */
    #earthvid {
        position: fixed;
        left: 50%;
        bottom: 0;
        transform: translateX(-50%);
        width: 100vw;
        min-width: 100vw;
        max-width: 100vw;
        height: 28vh;
        object-fit: cover;
        z-index: 3;
        opacity: 0.85;
        pointer-events: none;
        filter: blur(0.2px) brightness(1.2) saturate(1.25);
    }
    /* Spotlight/beam from above */
    .space-spotlight {
        position: fixed;
        top: -15vw;
        left: 50vw;
        transform: translateX(-50%);
        width: 120vw;
        height: 60vw;
        pointer-events: none;
        z-index: 2;
        background: radial-gradient(ellipse at 50% 0%, rgba(200,200,255,0.24) 0%, rgba(85,74,255,0.09) 54%, rgba(24,24,47,0.0) 100%);
        filter: blur(20px) brightness(1.12);
        animation: spotlight-fade 7s ease-in-out infinite alternate;
    }
    @keyframes spotlight-fade {
        0% { opacity: 0.83; }
        100% { opacity: 0.33; }
    }
    /* Massive faint watermark text */
    .watermark {
        position: fixed;
        top: 12vh;
        left: 50vw;
        transform: translateX(-50%);
        font-size: 13vw;
        font-family: 'Inter', sans-serif;
        color: #fff;
        opacity: 0.035;
        z-index: 4;
        letter-spacing: 0.04em;
        font-weight: 900;
        user-select: none;
        pointer-events: none;
        text-shadow: 0 0 80px #fff, 0 2px 16px #fff;
        white-space: nowrap;
    }
    /* Space particles with movement */
    .space-particles {
        position: fixed;
        top: 0; left: 0; width: 100vw; height: 100vh;
        pointer-events: none;
        z-index: 5;
    }
    .space-particles span {
        position: absolute;
        border-radius: 50%;
        background: #fff;
        opacity: 0.7;
        animation: move-particle 14s linear infinite, twinkle 2.5s infinite alternate;
    }
    /* Unique movement for each particle */
    .p1 { animation-delay: 0s, 0.2s; top:8%; left:15%; width:2.5px; height:2.5px; }
    .p2 { animation-delay: 2s, 0.7s; top:20%; left:60%; width:1.7px; height:1.7px; }
    .p3 { animation-delay: 3s, 1.1s; top:32%; left:80%; width:2.9px; height:2.9px; }
    .p4 { animation-delay: 1s, 0.9s; top:55%; left:30%; width:1.2px; height:1.2px; }
    .p5 { animation-delay: 6s, 0.4s; top:70%; left:10%; width:1.6px; height:1.6px; }
    .p6 { animation-delay: 8s, 1.5s; top:80%; left:85%; width:2.7px; height:2.7px; }
    .p7 { animation-delay: 7s, 0.6s; top:60%; left:50%; width:1.8px; height:1.8px; }
    .p8 { animation-delay: 5s, 1.3s; top:15%; left:78%; width:1.3px; height:1.3px; }
    .p9 { animation-delay: 4s, 0.8s; top:45%; left:65%; width:2.1px; height:2.1px; }
    .p10 { animation-delay: 10s, 1.1s; top:88%; left:40%; width:1.5px; height:1.5px; }
    .p11 { animation-delay: 12s, 1.6s; top:25%; left:35%; width:2.2px; height:2.2px; }
    .p12 { animation-delay: 13s, 1.7s; top:60%; left:80%; width:1.8px; height:1.8px; }
    .p13 { animation-delay: 11s, 1.2s; top:75%; left:55%; width:2.3px; height:2.3px; }
    .p14 { animation-delay: 9s, 0.5s; top:40%; left:20%; width:2.1px; height:2.1px; }
    .p15 { animation-delay: 7s, 1.8s; top:65%; left:70%; width:1.9px; height:1.9px; }
    .p16 { animation-delay: 3.5s, 1.4s; top:10%; left:90%; width:2.4px; height:2.4px; }
    .p17 { animation-delay: 2.5s, 0.3s; top:85%; left:20%; width:1.4px; height:1.4px; }
    .p18 { animation-delay: 4.5s, 1.9s; top:35%; left:55%; width:2.6px; height:2.6px; }
    .p19 { animation-delay: 6.5s, 1.1s; top:67%; left:33%; width:1.7px; height:1.7px; }
    .p20 { animation-delay: 8.5s, 0.6s; top:22%; left:77%; width:2.0px; height:2.0px; }
    @keyframes twinkle {
        from { opacity: 0.18; }
        to { opacity: 1; }
    }
    @keyframes move-particle {
        0% { transform: translateY(0) scale(1); }
        100% { transform: translateY(-30vh) scale(1.2); }
    }
    /* Glassmorphism for main container */
    .stApp {
        background: transparent !important;
    }
    .block-container {
        background: rgba(24, 24, 32, 0.62);
        border-radius: 20px;
        box-shadow: 0 8px 48px 0 rgba(80,40,200,0.18);
        padding: 2rem 2rem 2rem 2rem;
        backdrop-filter: blur(24px) saturate(1.7);
        border: 1.5px solid rgba(186,150,255,0.09);
    }
    /* Fading/gradient text effect for headings */
    h1, h2, h3, h4, h5, h6 {
        color: #fff;
        background: linear-gradient(90deg, #fff 0%, #b68cff 60%, #fff 100%);
        background-size: 200% auto;
        background-clip: text;
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        animation: fadeText 4s ease-in-out infinite alternate;
        font-weight: 900;
        letter-spacing: 0.01em;
        text-shadow: 0 2px 32px #a883ff33, 0 1px 2px #0008;
    }
    @keyframes fadeText {
        0% { background-position: 0% 50%; }
        100% { background-position: 100% 50%; }
    }
    .stButton>button {
        background: linear-gradient(90deg,#b68cff,#5f5fff);
        color: #fff;
        font-weight: 600;
        border-radius: 10px;

        box-shadow: 0 2px 12px rgba(120,80,255,0.13);
        transition: background 0.3s, box-shadow 0.3s;
    }
    .stButton>button:hover {
        background: linear-gradient(90deg,#5f5fff,#b68cff);
        box-shadow: 0 4px 24px rgba(186,150,255,0.22);
    }
    .stTextInput>div>input,
    .stTextArea>div>textarea {
        background: rgba(255,255,255,0.13);
        color: #fff;
        border-radius: 8px;
        border: 1px solid #b68cff;
        font-size: 1rem;
    }
    .stTabs [data-baseweb="tab"] {
        background: rgba(186,150,255,0.11);
        border-radius: 10px 10px 0 0;
        color: #fff;
        font-weight: 600;
    }
    .stTabs [aria-selected="true"] {
        background: #5f5fff;
        color: #fff;
    }
    </style>
    <video autoplay loop muted id="bgvid">
      <source src="https://assets.mixkit.co/videos/preview/mixkit-deep-space-stars-1600-large.mp4" type="video/mp4">
    </video>
    <video autoplay loop muted id="earthvid">
      <source src="https://cdn.pixabay.com/video/2023/05/28/163129-827828085_large.mp4" type="video/mp4">
    </video>
    <div class="space-spotlight"></div>
    <div class="watermark">AI Chat</div>
    <div class="space-particles">
      <span class="p1"></span><span class="p2"></span><span class="p3"></span><span class="p4"></span><span class="p5"></span>
      <span class="p6"></span><span class="p7"></span><span class="p8"></span><span class="p9"></span><span class="p10"></span>
      <span class="p11"></span><span class="p12"></span><span class="p13"></span><span class="p14"></span><span class="p15"></span>
      <span class="p16"></span><span class="p17"></span><span class="p18"></span><span class="p19"></span><span class="p20"></span>
    </div>
    ''',
    unsafe_allow_html=True
)
st.title("🤖 Ultimate AI Chatbot")

# ------------------- Emotion Detection ------------------- #
def record_audio(duration=15, samplerate=16000):
    st.info("Recording... Speak now!")
    audio = sd.rec(int(duration * samplerate), samplerate=samplerate, channels=1, dtype='int16')
    sd.wait()
    with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as temp_file:
        write(temp_file.name, samplerate, audio)
        return temp_file.name

def extract_basic_features(audio_data):
    energy = np.sum(audio_data.astype(np.float32)**2) / len(audio_data)
    zero_crossings = np.mean(np.abs(np.diff(np.sign(audio_data.flatten()))))
    max_amp = np.max(np.abs(audio_data))
    return energy, zero_crossings, max_amp

def detect_emotion_simple(audio_file):
    sr, data = read(audio_file)
    if len(data.shape) > 1:
        data = data[:, 0]
    energy, zcr, amp = extract_basic_features(data)
    if energy < 1000 and zcr < 0.05:
        return "Sad 😢"
    elif amp > 20000 and zcr > 0.1:
        return "Excited 🤩"
    elif energy < 5000:
        return "Calm 😌"
    else:
        return "Happy 😊"

def translate_text(text, target_lang):
    try:
        return GoogleTranslator(source='auto', target=target_lang).translate(text)
    except Exception as e:
        return f"[Translation Error] {e}"

def speak_translated_text(text, lang="en"):
    try:
        tts = gTTS(text=text, lang=lang)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as fp:
            tts.save(fp.name)
            st.audio(fp.name, format="audio/mp3")
    except Exception as e:
        st.error(f"Text-to-speech failed: {e}")

# -------------- Session State --------------
if "history" not in st.session_state: st.session_state.history = []
if "active_chat" not in st.session_state: st.session_state.active_chat = None
if "new_chat_mode" not in st.session_state: st.session_state.new_chat_mode = False
if "persona_prompt" not in st.session_state: st.session_state.persona_prompt = ""
if "selected_model" not in st.session_state: st.session_state.selected_model = "🧬 DeepSeek R1"

# ------------------- Custom Personas -------------------
chat_modes = {
    "🎓 Friendly Tutor": "You are a kind and patient tutor who explains concepts clearly, encourages curiosity, and keeps responses simple and friendly.",
    "👨‍💻 Strict Coder": "You are a strict senior developer. You give concise, direct answers, correct bad practices, and expect clean code.",
    "🧘 Motivational Coach": "You are a motivational coach who always responds positively, encourages the user, and uplifts their mindset.",
    "💼 Career Guide": "You are a career advisor helping users make decisions about their professional growth and opportunities.",
    "🤝 Friendly Chat": "You are a fun and casual chatbot. You use emojis and maintain a cheerful, light-hearted tone.",
    "🧠 Professional Assistant": "You are a knowledgeable professional. You give precise, formal answers suitable for business or technical discussions."
}

# -------------- Weather and News Functions --------------
def get_weather(city="Hyderabad"):
    url = f"http://api.openweathermap.org/data/2.5/weather?q={city}&appid={WEATHER_API_KEY}&units=metric"
    response = requests.get(url)
    if response.status_code == 200:
        data = response.json()
        return {"city": city, "temp": data["main"]["temp"], "description": data["weather"][0]["description"].capitalize()}
    return None

def get_news_from_rss(max_headlines=5):
    rss_url = "https://news.google.com/rss?hl=en-IN&gl=IN&ceid=IN:en"
    feed = feedparser.parse(rss_url)
    headlines = []
    for entry in feed.entries[:max_headlines]:
        headlines.append({"title": entry.title, "link": entry.link})
    return headlines

# -------------- Code Sandbox Functions -----------      
def execute_python_code(code, timeout=10):
    restricted_builtins = {
        'print': print, 'len': len, 'range': range, 'str': str, 'int': int, 'float': float,
        'list': list, 'dict': dict, 'tuple': tuple, 'set': set, 'bool': bool,
        'abs': abs, 'max': max, 'min': min, 'sum': sum, 'sorted': sorted,
        'reversed': reversed, 'enumerate': enumerate, 'zip': zip,
        'map': map, 'filter': filter, 'all': all, 'any': any, 'round': round,
        'pow': pow, 'divmod': divmod, 'type': type, 'isinstance': isinstance,
        'hasattr': hasattr, 'getattr': getattr, 'setattr': setattr,
        'chr': chr, 'ord': ord, 'bin': bin, 'oct': oct, 'hex': hex
    }
    safe_modules = {'math', 'random', 'datetime', 'json', 're', 'collections',
                    'itertools', 'functools', 'operator', 'string', 'textwrap'}
    def safe_import(name, globals=None, locals=None, fromlist=(), level=0):
        if name in safe_modules:
            return __import__(name, globals, locals, fromlist, level)
        else:
            raise ImportError(f"Module '{name}' is not allowed in sandbox")
    sandbox_globals = {'__builtins__': restricted_builtins, '__import__': safe_import}
    output_buffer = io.StringIO()
    error_buffer = io.StringIO()
    result = {'success': False, 'output': '', 'errors': 'Unknown error occurred'}
    def target():
        try:
            with redirect_stdout(output_buffer), redirect_stderr(error_buffer):
                exec(code, sandbox_globals)
            result.update({'success': True, 'output': output_buffer.getvalue(), 'errors': error_buffer.getvalue() or None})
        except Exception as e:
            result.update({'output': output_buffer.getvalue(), 'errors': f"{type(e).__name__}: {str(e)}\n{traceback.format_exc()}"})
    thread = threading.Thread(target=target)
    thread.start()
    thread.join(timeout)
    if thread.is_alive():
        result.update({'success': False, 'output': '', 'errors': f'Code execution timed out (limit: {timeout}s)'})
        thread.join()
    return result

# -------------- Offline Mode Functions --------------
def check_internet_connection():
    try:
        response = requests.get('https://www.google.com',  timeout=5)
        return response.status_code == 200
    except:
        return False

def cache_data(data, cache_type, filename):
    cache_dir = os.path.join(SAVE_DIR, "cache", cache_type)
    os.makedirs(cache_dir, exist_ok=True)
    cache_file = os.path.join(cache_dir, f"{filename}.json")
    try:
        with open(cache_file, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        return True
    except:
        return False

def load_cached_data(cache_type, filename):
    cache_file = os.path.join(SAVE_DIR, "cache", cache_type, f"{filename}.json")
    try:
        if os.path.exists(cache_file):
            with open(cache_file, 'r', encoding='utf-8') as f:
                return json.load(f)
    except:
        pass
    return None

def get_offline_weather(city):
    cached = load_cached_data("weather", city.lower())
    return cached if cached else None

def get_offline_news():
    cached = load_cached_data("news", "headlines")
    return cached if cached else []

if "is_online" not in st.session_state:
    st.session_state.is_online = check_internet_connection()
if not st.session_state.is_online:
    st.warning("🔌 You're offline - using cached content and limited functionality")


# ========== ⭐ FUNCTION: Auto Persona Detection ==========

# --- Initialize session state variables ---
if "manual_persona_enabled" not in st.session_state:
    st.session_state.manual_persona_enabled = False

def detect_persona_from_input(user_input):
    input_lower = user_input.lower()
    if any(word in input_lower for word in ["code", "function", "bug", "python", "script", "error", "logic", "syntax"]):
        return chat_modes["👨‍💻 Strict Coder"]
    elif any(word in input_lower for word in ["explain", "learn", "study", "understand", "concept", "exam", "doubt"]):
        return chat_modes["🎓 Friendly Tutor"]
    elif any(word in input_lower for word in ["depressed", "motivate", "inspire", "confidence", "failure", "success", "emotion"]):
        return chat_modes["🧘 Motivational Coach"]
    elif any(word in input_lower for word in ["resume", "job", "career", "interview", "linkedin", "cv"]):
        return chat_modes["💼 Career Guide"]
    elif any(word in input_lower for word in ["fun", "lol", "joke", "casual", "friend", "hello", "hey"]):
        return chat_modes["🤝 Friendly Chat"]
    elif any(word in input_lower for word in ["project", "meeting", "report", "company", "client", "professional"]):
        return chat_modes["🧠 Professional Assistant"]
    else:
        return chat_modes["🤝 Friendly Chat"]  # default fallback


# Login Page
if "email" not in st.session_state:
    st.title("🔐 Login/Signup To Use ChatBOT")
    tab1, tab2, tab3 = st.tabs(["Sign Up", "Login", "Forgot Password"])
    with tab1:
        if "signup_email" not in st.session_state:
            new_email = st.text_input("Email")
            new_pass = st.text_input("Password", type="password")
            if new_pass:
                if is_strong_password(new_pass):
                    st.success("✅ Strong Password")
                else:
                    st.warning("❌ Weak password. Use 8+ chars with upper, lower, number & symbol.")
            if st.button("Register"):
                if not validate_email(new_email):
                    st.error("Invalid email format.")
                elif not is_strong_password(new_pass):
                    st.error("Password too weak.")
                else:
                    hashed = hash_password(new_pass)
                    code = str(uuid.uuid4())[:6]
                    try:
                        c.execute("INSERT INTO users VALUES (?, ?, ?, ?)", (new_email, hashed, 0, code))
                        conn.commit()
                        send_verification_email(new_email, code)
                        st.success("✅ Registered. Check your email for the OTP.")
                        st.session_state.signup_email = new_email
                    except sqlite3.IntegrityError:
                        st.warning("Email already exists. Try logging in.")
        else:
            otp_signup = st.text_input("Enter OTP sent to your email")
            if st.button("Verify OTP"):
                c.execute("SELECT verification_code FROM users WHERE email=?", (st.session_state.signup_email,))
                code = c.fetchone()[0]
                if otp_signup == code:
                    c.execute("UPDATE users SET is_verified=1 WHERE email=?", (st.session_state.signup_email,))
                    conn.commit()
                    st.success("Email verified ✅ You can now login.")
                    del st.session_state.signup_email
                    st.rerun()
                else:
                    st.error("Invalid OTP")
    with tab2:
        login_email = st.text_input("Email", key="login")
        login_pass = st.text_input("Password", type="password", key="login_pw")
        if st.button("Login"):
            c.execute("SELECT password, is_verified FROM users WHERE email=?", (login_email,))
            user = c.fetchone()
            if not user:
                st.error("No such user.")
            elif user[0] != hash_password(login_pass):
                st.error("Incorrect password.")
            elif not user[1]:
                st.warning("Verify your email via OTP.")
            else:
                st.session_state.email = login_email
                st.success("Logged in ✅")
                st.rerun()
    with tab3:
        forgot_email = st.text_input("Enter your registered email")
        if st.button("Send OTP"):
            c.execute("SELECT * FROM users WHERE email=?", (forgot_email,))
            if not c.fetchone():
                st.error("Email not found.")
            else:
                reset_code = str(uuid.uuid4())[:6]
                c.execute("UPDATE users SET verification_code=? WHERE email=?", (reset_code, forgot_email))
                conn.commit()
                send_verification_email(forgot_email, reset_code)
                st.session_state.reset_email = forgot_email
                st.session_state.otp_sent = True
                st.success("OTP sent to your email.")
    if st.session_state.get("otp_sent"):
        st.subheader("🔁 Reset Password")
        entered_otp = st.text_input("Enter OTP")
        new_reset_pass = st.text_input("New Password", type="password")
        if new_reset_pass and not is_strong_password(new_reset_pass):
            st.warning("❌ Weak password. Use 8+ chars with upper, lower, number & symbol.")
        if st.button("Reset Password"):
            c.execute("SELECT verification_code FROM users WHERE email=?", (st.session_state.reset_email,))
            stored_code = c.fetchone()
            if stored_code and entered_otp == stored_code[0]:
                if is_strong_password(new_reset_pass):
                    hashed_new = hash_password(new_reset_pass)
                    c.execute("UPDATE users SET password=?, verification_code=? WHERE email=?",
                              (hashed_new, "", st.session_state.reset_email))
                    conn.commit()
                    st.success("Password reset successful ✅")
                    st.session_state.otp_sent = False
                    st.session_state.reset_email = ""
                    st.rerun()
                else:
                    st.error("Password too weak.")
            else:
                st.error("Invalid OTP")
    st.stop()

# After login, verify OTP if needed
c.execute("SELECT is_verified FROM users WHERE email=?", (st.session_state.email,))
verified = c.fetchone()[0]
if not verified:
    st.info("Enter OTP sent to your email")
    otp_input = st.text_input("OTP")
    if st.button("Verify"):
        c.execute("SELECT verification_code FROM users WHERE email=?", (st.session_state.email,))
        code = c.fetchone()[0]
        if otp_input == code:
            c.execute("UPDATE users SET is_verified=1 WHERE email=?", (st.session_state.email,))
            conn.commit()
            st.success("Email verified ✅")
            st.rerun()
        else:
            st.error("Invalid OTP")
    st.stop()

# Main Chat UI starts here
st.markdown("---")
username_only = st.session_state.email.split("@")[0]
st.header(f"Welcome, {username_only}")

if "history" not in st.session_state:
    st.session_state.history = []
    c.execute("SELECT user_msg, bot_msg, timestamp FROM chat_history WHERE email=? ORDER BY timestamp", 
              (st.session_state.email,))
    rows = c.fetchall()
    for row in rows:
        st.session_state.history.append((row[0], row[1]))

def save_chat_to_db(user_input, response):
    timestamp = datetime.now().isoformat()
    chat_id = str(uuid.uuid4())
    c.execute("INSERT INTO chat_history VALUES (?, ?, ?, ?, ?)", 
              (chat_id, st.session_state.email, user_input, response, timestamp))
    conn.commit()

def web_search(query, num=5):
    url = "https://google.serper.dev/search"   
    SERPER_API_KEY = "be079c48e1e1df5928813172860bc4d5e73b6a16"  # 🔑 Paste key here

    headers = {
        "X-API-KEY": SERPER_API_KEY,
        "Content-Type": "application/json"
    }
    payload = {"q": query}
    try:
        response = requests.post(url, headers=headers, json=payload, timeout=10)
        results = response.json().get("organic", [])
        return [f"🔗 {r['link']}\n📄 {r['title']}" for r in results[:num]]
    except Exception as e:
        return [f"❌ Web search failed: {e}"]

# ========== ⭐ SIDEBAR SETTINGS (Updated) ==========
with st.sidebar:
    st.header("💾 Chat Management")
    if st.button("🆕 New Chat"):
        st.session_state.history = []
        st.session_state.active_chat = None
        st.session_state.new_chat_mode = True
        st.rerun()

    files = sorted([f for f in os.listdir(SAVE_DIR) if f.endswith(".json")], reverse=True)
    q = st.text_input("🔍 Search Saved Chats")
    sel = st.selectbox("📂 Load Chat", ["--"] + [f for f in files if q.lower() in f.lower()])
    if sel != "--" and sel:
        st.session_state.history = json.load(open(os.path.join(SAVE_DIR, sel), encoding="utf-8"))
        st.session_state.active_chat = sel
        st.session_state.new_chat_mode = False
        st.success(f"✅ Loaded: {sel}")
    if st.session_state.active_chat:
        nm = st.text_input("✏️ Rename Chat", value=st.session_state.active_chat.split("__")[-1].replace(".json", ""))
        if st.button(".btnSave"):
            ts = st.session_state.active_chat.split("__")[0]
            new_fname = f"{ts}__{'_'.join(nm.lower().split())}.json"
            os.rename(os.path.join(SAVE_DIR, st.session_state.active_chat), os.path.join(SAVE_DIR, new_fname))
            st.session_state.active_chat = new_fname
            st.success(f"✅ Renamed to: {new_fname}")

    st.markdown("---")
    if st.button("⚙️ Settings", key="sidebar_settings_btn"):
        st.session_state.show_settings = True

# ===== Settings Modal/Expandable =====
if st.session_state.get("show_settings"):
    with st.container():
        st.markdown("""
            <div style='position:fixed; top:5vh; left:50vw; transform:translateX(-50%); z-index:10000; background:rgba(34,34,54,0.97); border-radius:18px; box-shadow:0 4px 32px #b68cff44; padding:2.5rem 2.5rem 2rem 2.5rem; min-width:340px; max-width:90vw;'>
        """, unsafe_allow_html=True)
        st.markdown("### ⚙️ Settings")
        if st.button("❌ Close Settings", key="close_settings_btn"):
            st.session_state.show_settings = False
            st.rerun()
        st.markdown("---")
        # Persona selection
        st.markdown("#### 🧠 Chat Persona (Auto Detection)")
        st.session_state.manual_persona_enabled = st.checkbox("Enable Manual Persona Selection", value=st.session_state.get("manual_persona_enabled", False), key="modal_manual_persona")
        if st.session_state.manual_persona_enabled:
            selected_mode = st.selectbox("Select Chat Persona", list(chat_modes.keys()), key="modal_select_persona")
            st.session_state.persona_prompt = chat_modes[selected_mode]
            st.markdown(f"**Manual Persona Set:** {selected_mode}")
        else:
            st.markdown("🔍 Running in **Auto Persona Detection** mode")
        st.markdown("---")
        # Model selection
        st.markdown("#### 🤖 Model Selection")
        st.session_state.selected_model = st.selectbox(
            "Choose LLM Model", list(model_sources.keys()), key="modal_select_model"
        )
        st.markdown("</div>", unsafe_allow_html=True)

# ===== Tabs =====
tabs = st.tabs([
    "💬 Chat", "🎨 Image", "🪄 PPT/PDF", "📺 YouTube Summary", 
    "📄 Resume Review", "💻 Code Tools", "📅 Daily Utilities", 
    "🧪 Code Sandbox", "Offline Mode", "🎙️ Voice Features", 
    "🛠️ AI Tools", "🎮 Game Center", "🩺 Health Assistant"
])

# ================= ENHANCED FEATURES ===================
# --- Theme Toggle ---
theme = st.session_state.get("theme", "auto")
theme = st.selectbox("Theme", ["auto", "light", "dark"], index=["auto", "light", "dark"].index(theme), key="theme_toggle")
st.session_state.theme = theme
st.markdown(f"<style>body {{ background: {'#232323' if theme=='dark' else '#fff'}; }}</style>", unsafe_allow_html=True)

# --- Collapsible Chat History ---
show_history = st.session_state.get("show_history", True)
if st.button("📜 Toggle Chat History", key="toggle_history_btn"):
    st.session_state.show_history = not show_history
show_history = st.session_state.get("show_history", True)

# --- Profile/Settings Modal ---
with st.expander("👤 Profile & Settings", expanded=False):
    st.write(f"**User:** {st.session_state.get('email', 'Guest')}")
    st.write("**Achievements:**", st.session_state.get("achievements", []))
    st.write("**Theme:**", st.session_state.get("theme", "auto"))
    st.write("**Chats:**", len(st.session_state.get("history", [])))
    st.write("(More settings coming soon!)")

# --- Daily AI Challenge/Quiz ---
with st.expander("🎯 Daily AI Challenge", expanded=False):
    quiz = st.session_state.get("daily_quiz", {
        "q": "What is the capital of France?",
        "a": "Paris"
    })
    if not isinstance(quiz, dict):
        quiz = {
            "q": "What is the capital of France?",
            "a": "Paris"
        }
    st.write(f"**Today's Quiz:** {quiz['q']}")
    user_ans = st.text_input("Your Answer", key="daily_quiz")
    if user_ans and user_ans.lower().strip() == quiz["a"].lower():
        st.success("Correct! 🏆")
        st.session_state.achievements = st.session_state.get("achievements", []) + ["Daily Quiz Winner"]
    elif user_ans:
        st.warning("Try again!")

# --- Chat Analytics Dashboard ---
with st.expander("📊 Chat Analytics", expanded=False):
    hist = st.session_state.get("history", [])
    st.write(f"**Total Messages:** {len(hist)}")
    st.write(f"**Unique Days Active:** {len(set([str(m[2])[:10] if len(m)>2 else 'today' for m in hist]))}")
    st.write(f"**Achievements:** {st.session_state.get('achievements', [])}")
    # (Extend with more analytics as needed)

# --- Chat History with Avatars, Reactions, Voice, Quick Actions, Feedback ---
if show_history:
    st.markdown('<div class="chat-container">', unsafe_allow_html=True)
    for i, (q, a, *rest) in enumerate(st.session_state.get("history", [])):
        st.markdown(f'<div class="msg-right"><img src="https://api.dicebear.com/7.x/personas/svg?seed=user" width="32" style="vertical-align:middle;border-radius:50%;margin-right:8px;"> {q}</div>', unsafe_allow_html=True)
        st.markdown(f'<div class="msg-left"><img src="https://api.dicebear.com/7.x/bottts/svg?seed=bot" width="32" style="vertical-align:middle;border-radius:50%;margin-right:8px;"> {a}', unsafe_allow_html=True)
        # Emoji reactions (safe session_state handling)
        like_key = f"like_{i}"
        dislike_key = f"dislike_{i}"
        if like_key not in st.session_state:
            st.session_state[like_key] = 0
        if dislike_key not in st.session_state:
            st.session_state[dislike_key] = 0
        if st.button("👍", key=like_key):
            st.session_state[like_key] += 1
        if st.button("👎", key=dislike_key):
            st.session_state[dislike_key] += 1
        st.write(f"Likes: {st.session_state.get(like_key, 0)} | Dislikes: {st.session_state.get(dislike_key, 0)}")
        # Voice output
        if st.button("🔊 Listen", key=f"tts_{i}"):
            tts = gTTS(a)
            tts.save("bot_voice.mp3")
            audio_file = open("bot_voice.mp3", "rb")
            st.audio(audio_file.read(), format="audio/mp3")
        # Quick actions
        st.write("Quick Actions:")
        if st.button("Summarize", key=f"summarize_{i}"): st.info("(Stub) Summary: " + a[:50] + "...")
        if st.button("Translate", key=f"translate_{i}"): st.info("(Stub) Translation: " + a)
        if st.button("Explain", key=f"explain_{i}"): st.info("(Stub) Explanation: " + a)
        # Feedback
        st.write("Rate this response:")
        rating = st.slider("", 1, 5, 3, key=f"rating_{i}")
        st.session_state[f"rating_{i}"] = rating
        st.write(f"Your rating: {rating}")
        st.markdown('</div>', unsafe_allow_html=True)
    st.markdown('<div id="scroll-bottom" class="scroll-bottom"></div>', unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)
    st.markdown('''<script>var chatBottom = document.getElementById('scroll-bottom');if (chatBottom) { chatBottom.scrollIntoView({behavior: "smooth"}); }</script>''', unsafe_allow_html=True)

# --- Bot Typing Indicator (stub, to be triggered on response fetch) ---
if st.session_state.get("bot_typing", False):
    st.info("🤖 Bot is typing...")

# --- Accessibility Improvements ---
st.markdown("<style>:focus {outline: 2px solid #a259ff !important;} .chat-container {font-size: 1.1rem;} .msg-right,.msg-left {padding: 8px 16px;margin: 8px 0;border-radius: 14px;max-width: 70%;display: inline-block;} .msg-right {background: #a259ff22;float: right;clear: both;} .msg-left {background: #23232322;float: left;clear: both;} </style>", unsafe_allow_html=True)

# --- Stubs for future features ---
# - Plugin system, PWA, encryption, API/webhook, calendar integration, E2E encryption, etc.
#   (See comments for future modularization.)
# ======================================================

# ========== 🅰️ CHAT TAB: Persona-Aware ==========
with tabs[0]:
    st.subheader("💬 Ask Anything")
    # --- Plus Button and Feature Menu ---
    col1, col2 = st.columns([10,1])
    with col1:
        pass  # Removed duplicate chat_input from col1
    with col2:
        if st.button("➕", key="plus_btn"):
            st.session_state.show_plus_menu = not st.session_state.get("show_plus_menu", False)
    user_input = st.chat_input("Type your message...", key="main_chat_input")
    # --- Ensure variables are always defined, even if plus menu is closed ---
    uploaded_file = None
    add_text = ""
    web_enabled = False

    if st.session_state.get("show_plus_menu", False):
        with st.expander("More Features", expanded=True):
            # Upload a file
            uploaded_file = st.file_uploader("📄 Upload PDF/DOCX/TXT", type=["pdf", "docx", "txt"], key="plus_file")
            # Add text content
            add_text = st.text_area("📝 Add Text Content", key="plus_text")
            # Enable web search
            web_enabled = st.checkbox("🌐 Enable Web Search", key="plus_web")
            # Voice input
            # Voice input (check if PyAudio is installed)
            try:
                import pyaudio
                pyaudio_available = True
            except ImportError:
                pyaudio_available = False
            if pyaudio_available:
                if st.button("🎤 Voice Input", key="plus_voice"):
                    r = sr.Recognizer()
                    with sr.Microphone() as src:
                        st.info("🎤 Listening...")
                        audio = r.listen(src)
                    try:
                        user_input = r.recognize_google(audio)  # type: ignore[attr-defined]
                        st.session_state["main_chat_input"] = user_input
                    except:
                        st.warning("Could not recognize your voice.")
                        user_input = ""
            else:
                st.info("PyAudio is not installed. Voice input is disabled. To use voice input, install PyAudio (pip install pyaudio).")

    # 💡 Handle message from user
    if user_input:
        with st.spinner("🤖 Thinking..."):
            webctx, docctx = "", ""

            # Web search context
            if web_enabled:
                web_results = web_search(user_input)
                webctx = "\n".join(web_results)

            # Document upload context
            docctx = ""
            if uploaded_file:
                try:
                    if uploaded_file.type.endswith("pdf"):
                        reader = PyPDF2.PdfReader(uploaded_file)
                        docctx = " ".join([p.extract_text() or "" for p in reader.pages])
                    elif uploaded_file.type.endswith("docx"):
                        doc = docx.Document(uploaded_file)
                        docctx = " ".join([p.text for p in doc.paragraphs])
                    elif uploaded_file.type.endswith("txt"):
                        uploaded_file.seek(0)
                        docctx = uploaded_file.read().decode("utf-8", errors="ignore")
                except Exception as e:
                    st.warning(f"❌ Could not read uploaded file: {e}")
                    docctx = ""
                if docctx.strip():
                    st.info(f"**File content preview:**\n{docctx[:500]}" + ("..." if len(docctx) > 500 else ""))
                else:
                    st.warning("⚠️ Uploaded file appears empty or could not be read. Please check your file.")

            # 🔍 Auto or Manual Persona Selection
            if not st.session_state.manual_persona_enabled:
                st.session_state.persona_prompt = detect_persona_from_input(user_input)

            # Full prompt construction
            full_prompt = f"{webctx}\n\n{docctx}\n\nUser: {user_input}"

            messages = [
                {"role": "system", "content": st.session_state.persona_prompt},
                {"role": "user", "content": full_prompt}
            ]

            response_text, model_used = get_ai_response(messages, model_preference=st.session_state.selected_model)
            tagged_res = f"**🤖 Response from {model_used}:**\n{response_text}"

            # Save and Display
            st.session_state.history.append((user_input, tagged_res))
            save_chat_to_db(user_input, tagged_res)

            if st.session_state.new_chat_mode or not st.session_state.active_chat:
                fname = f"{datetime.now().strftime('%Y%m%d%H%M%S')}__{'_'.join(user_input.split()[:3])}.json"
                with open(os.path.join(SAVE_DIR, fname), "w", encoding="utf-8") as f:
                    json.dump(st.session_state.history, f, ensure_ascii=False, indent=2)
                st.session_state.active_chat = fname
                st.session_state.new_chat_mode = False

    # --- Custom CSS for left/right alignment ---
    st.markdown("""
    <style>
    .stApp {
        background: #181824 !important;
        color: #fff !important;
    }
    .chat-container {
        background: transparent !important;
        max-height: 500px;
        overflow-y: auto;
        padding-bottom: 16px;
    }
    .msg-left {
        background: #fff !important;
        color: #232323 !important;
        border-radius: 16px 16px 6px 16px;
        padding: 13px 18px;
        margin: 10px 0;
        max-width: 70%;
        float: left;
        clear: both;
        box-shadow: 0 2px 12px 0 #00000020;
        font-size: 1.08rem;
        border: none !important;
    }
    .msg-right {
        background: #fff !important;
        color: #232323 !important;
        border-radius: 16px 16px 6px 16px;
        padding: 13px 18px;
        margin: 10px 0;
        max-width: 70%;
        float: right;
        clear: both;
        text-align: right;
        box-shadow: 0 2px 12px 0 #00000020;
        font-size: 1.08rem;
        border: none !important;
    }
    .stChatInputContainer, .chat-input {
        background: #fff !important;
        border-radius: 20px !important;
        box-shadow: 0 2px 16px 0 #00000044;
        padding: 10px 18px !important;
        margin-top: 10px;
        margin-bottom: 0;
        width: 100% !important;
        border: none !important;
        color: #232323 !important;
    }
    .stChatInputContainer input, .chat-input input {
        background: transparent !important;
        color: #232323 !important;
        font-size: 1.1rem;
        flex: 1;
        outline: none;
    }
    .stChatInputContainer button, .chat-input button {
        background: #232323 !important;
        color: #fff !important;
        border-radius: 50% !important;
        box-shadow: 0 0 8px 1px #00000033 !important;
        width: 40px !important;
        height: 40px !important;
        margin-left: 10px !important;
        font-size: 1.3rem !important;
        cursor: pointer;
        border: none !important;
        transition: box-shadow 0.2s, background 0.2s;
    }
    .stChatInputContainer button:hover, .chat-input button:hover {
        box-shadow: 0 0 16px 4px #23232366 !important;
        background: #444 !important;
    }
    .stTabs [data-baseweb="tab-list"] {
        background: #232323 !important;
        border-radius: 20px !important;
        box-shadow: 0 2px 12px 0 #00000020;
        border: none !important;
    }
    .stTabs [data-baseweb="tab"] {
        color: #fff !important;
        border: none !important;
        background: transparent !important;
    }
    .stTabs [aria-selected="true"] {
        background: #fff !important;
        color: #232323 !important;
        border: none !important;
    }
    .scroll-bottom { height: 1px; }
    ::selection { background: #23232322; }
    </style>
""", unsafe_allow_html=True)


    # --- Chat Messages ---
    st.markdown('<div class="chat-container">', unsafe_allow_html=True)
    for q, a in st.session_state.history:
        st.markdown(f'<div class="msg-right">{q}</div>', unsafe_allow_html=True)
        st.markdown(f'<div class="msg-left">{a}</div>', unsafe_allow_html=True)
    st.markdown('<div id="scroll-bottom" class="scroll-bottom"></div>', unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)

    # --- Auto-scroll to bottom using JS ---
    st.markdown('''
        <script>
        var chatBottom = document.getElementById('scroll-bottom');
        if (chatBottom) { chatBottom.scrollIntoView({behavior: "smooth"}); }
        </script>
    ''', unsafe_allow_html=True)


# PPT / PDF Generator Tab
with tabs[2]:
    st.subheader("🪄 Generate PPT / PDF from Prompt")
    THEME_COLORS = [
        RGBColor(240, 240, 255),
        RGBColor(255, 240, 230),
        RGBColor(230, 245, 255),
        RGBColor(255, 240, 230),
        RGBColor(235, 255, 235),
        RGBColor(255, 245, 230)
    ]
    def generate_slides(topic, num_slides=5):
        messages = [{
            "role": "user",
            "content": f"""
            Generate a detailed PowerPoint presentation with {num_slides} slides on the topic: "{topic}".
            Each slide should include:
            - A title
            - 5 or more bullet points with clear explanations
            - An image description for visual aid
            Keep each bullet under 20 words.
            Format the response as a JSON list like this:
            [
              {{
                "title": "...",
                "bullets": ["...", "...", "..."],
                "image_prompt": "..."
              }},
              ...
            ]
            """
        }]
        response_text, model_used = get_ai_response(messages)
        try:
            first_bracket = response_text.find("[")
            last_bracket = response_text.rfind("]")
            clean_json = response_text[first_bracket:last_bracket+1]
            return json.loads(clean_json)
        except Exception as e:
            return []
    def create_ppt(slides):
        prs = Presentation()
        for slide_data in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = random.choice(THEME_COLORS)
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = slide_data.get("title", "Untitled")
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            bullets = slide_data.get("bullets", [])
            bullet_count = len(bullets)
            bullet_box_height = min(0.5 * bullet_count, 3.8)
            content_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3), 
                Inches(5.5), Inches(bullet_box_height)
            )
            tf = content_shape.text_frame
            tf.clear()
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
            image_prompt = slide_data.get("image_prompt")
            if image_prompt:
                try:
                    url = f"https://image.pollinations.ai/prompt/{urllib.parse.quote(image_prompt)}"  
                    response = requests.get(url)
                    img = Image.open(BytesIO(response.content))
                    img.thumbnail((400, 400))
                    img_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                    img.save(img_file.name)
                    if bullet_box_height < 2.5:
                        img_left = Inches(2.5)
                        img_top = Inches(1.3 + bullet_box_height + 0.3)
                    else:
                        img_left = Inches(6.1)
                        img_top = Inches(1.6)
                    slide.shapes.add_picture(img_file.name, img_left, img_top, width=Inches(3.0))
                except Exception as e:
                    st.warning(f"⚠️ Could not generate image: {str(e)}")
        return prs
    def create_pdf(slides):
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Arial", size=14)
        for i, slide_data in enumerate(slides, 1):
            pdf.cell(0, 10, f"Slide {i}: {slide_data.get('title', 'Untitled')}", ln=True, align='L')
            pdf.set_font("Arial", size=12)
            bullets = slide_data.get("bullets", [])
            for bullet in bullets:
                pdf.cell(0, 10, f"• {bullet}", ln=True, align='L')
            pdf.ln(10)
            image_prompt = slide_data.get("image_prompt")
            if image_prompt:
                try:
                    url = f"https://image.pollinations.ai/prompt/{urllib.parse.quote(image_prompt)}"
                    response = requests.get(url)
                    img = Image.open(BytesIO(response.content))
                    img_path = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                    img.save(img_path.name)
                    # Always insert image below text, with at least 8 units gap
                    y_now = pdf.get_y()
                    img_height = 60  # Approximate, will scale
                    if y_now + img_height + 15 > 280:
                        pdf.add_page()
                        y_now = pdf.get_y()
                    pdf.ln(8)
                    pdf.image(img_path.name, x=10, y=y_now+8, w=180)
                    pdf.ln(img_height+8)
                except Exception:
                    pdf.cell(0, 10, "[Image could not be generated]", ln=True)
            pdf.add_page()
        return pdf
    topic = st.text_area("Enter topic or prompt:")
    col1, col2 = st.columns([3, 1])
    with col1:
        num_slides = st.slider("Select number of slides", 3, 20, 5)
    with col2:
        file_type = st.selectbox("Choose file type:", ["pptx", "pdf"], key="ppt_pdf_type")
    if st.button("Generate Slides") and topic:
        with st.spinner("AI is creating your presentation..."):
            slides = generate_slides(topic, num_slides)
            if slides:
                if file_type == "pptx":
                    doc = create_ppt(slides)
                    file_ext = ".pptx"
                    file_name = f"{topic[:20].strip().replace(' ', '_')}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx"
                    doc.save(file_name)
                else:
                    doc = create_pdf(slides)
                    file_ext = ".pdf"
                    file_name = f"{topic[:20].strip().replace(' ', '_')}_{datetime.now().strftime('%Y%m%d%H%M%S')}.pdf"
                    doc.output(file_name)
                with open(file_name, "rb") as f:
                    b64 = base64.b64encode(f.read()).decode()
                    href = f'<a href="data:application/octet-stream;base64,{b64}" download="{file_name}">📥 Download {file_ext.upper()}</a>'
                    st.markdown(href, unsafe_allow_html=True)
            else:
                st.error("Failed to generate slides. Please try again.")

# YouTube Summarization
with tabs[3]:
    import urllib.parse
    import langdetect
    from gtts import gTTS
    import tempfile
    st.subheader("📺 YouTube Video Summary")
    yt_url = st.text_input("Paste YouTube URL:")
    preferred_output_lang = st.radio(
        "Show Summary In:",
        ["Same as Video Language", "English Translation"],
        index=1
    )
    def extract_video_id(url):
        try:
            parsed = urllib.parse.urlparse(url)
            if parsed.netloc == "www.youtube.com":
                return urllib.parse.parse_qs(parsed.query).get("v", [None])[0]
            elif parsed.netloc == "youtu.be":
                return parsed.path[1:]
        except:
            return None
    def detect_language(text):
        from langdetect import detect
        try:
            return detect(text)
        except:
            return "en"
    if st.button("Summarize Video") and yt_url:
        try:
            from youtube_transcript_api._api import YouTubeTranscriptApi
            from youtube_transcript_api._errors import TranscriptsDisabled, NoTranscriptFound
            vid_id = extract_video_id(yt_url)
            if not vid_id:
                st.error("❌ Could not extract video ID.")
            else:
                try:
                    transcript_list = YouTubeTranscriptApi.list_transcripts(vid_id)
                    try:
                        transcript = transcript_list.find_transcript(["en"])
                    except:
                        transcript = transcript_list.find_transcript(["hi", "te", "ta", "ml", "bn"])
                    fetched = transcript.fetch()
                    full_text = " ".join([t.text for t in fetched])
                    detected_lang = detect_language(full_text)
                    messages = [{
                        "role": "user",
                        "content": f"Summarize this in short points:\n{full_text}"
                    }]
                    summary, model_used = get_ai_response(messages)
                    if preferred_output_lang == "English Translation" and detected_lang != "en":
                        summary = translate_text(summary, "en")
                    st.success("✅ Summary:")
                    st.markdown(summary)
                    audio_lang = 'en' if preferred_output_lang == "English Translation" else detected_lang
                    tts = gTTS(text=summary, lang=audio_lang if audio_lang in ['en', 'hi', 'te'] else 'en')
                    audio_path = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3").name
                    tts.save(audio_path)
                    st.audio(audio_path, format="audio/mp3")
                except (TranscriptsDisabled, NoTranscriptFound):
                    st.error("❌ No subtitles found for this video in any language.")
        except Exception as e:
            st.error(f"❌ Failed to process: {e}")

# Resume Review Tab
with tabs[4]:
    st.subheader("📄 Resume Analyzer")
    resume = st.file_uploader("Upload Resume (PDF or DOCX)", type=["pdf", "docx"], key="resume_upload")
    job = st.text_area("Paste job description:")
    if st.button("Analyze Resume") and resume and job:
        text = ""
        if resume.type.endswith("pdf"):
            reader = PyPDF2.PdfReader(resume)
            text = "\n".join([p.extract_text() or "" for p in reader.pages])
        elif resume.type.endswith("docx"):
            doc = docx.Document(resume)
            text = "\n".join([p.text for p in doc.paragraphs])
        prompt = f"Evaluate this resume:\n{text}\nFor this job:\n{job}\nGive actionable feedback and score out of 100."
        messages = [{"role": "user", "content": prompt}]
        response_text, model_used = get_ai_response(messages)
        st.markdown(f"**🤖 Response from {model_used}:**\n\n{response_text}")

# Code Assistant Tab
with tabs[5]:
    st.subheader("💻 Code Assistant")
    tool = st.selectbox("Choose Tool", ["Explain", "Refactor", "Generate"])
    code_input = st.text_area("Paste your code or requirement:")
    if st.button("Run Code Tool") and code_input:
        system_prompt = {
            "Explain": f"Explain this code clearly:\n{code_input}",
            "Refactor": f"Refactor this code and explain the changes:\n{code_input}",
            "Generate": f"Generate code for this:\n{code_input}"
        }
        messages = [{"role": "user", "content": system_prompt[tool]}]
        response_text, model_used = get_ai_response(messages)
        st.markdown(f"**🤖 Response from {model_used}:**\n\n{response_text}")

# Daily Utilities Tab
with tabs[6]:
    st.title("📅 Daily Utilities")
    today = datetime.now()
    st.markdown(f"### 📅 Today: `{today.strftime('%A, %d %B %Y')}`")
    city = st.text_input("Enter city for weather", "Hyderabad")
    if st.session_state.is_online:
        weather = get_weather(city)
        if weather:
            st.success(f"🌦️ Weather in **{weather['city']}**: {weather['temp']}°C, {weather['description']}")
            cache_data(weather, "weather", city.lower())
        else:
            st.error("❌ Weather data unavailable")
    else:
        cached_weather = get_offline_weather(city)
        if cached_weather:
            st.info(f"🌦️ Cached Weather in **{cached_weather['city']}**: {cached_weather['temp']}°C, {cached_weather['description']}")
        else:
            st.warning("❌ No cached weather data available")
    st.markdown("### 🗞️ Top Headlines (Google News)")
    # Date-based cache for daily headlines
    today_str = today.strftime('%Y-%m-%d')
    cached_date = st.session_state.get('news_cache_date')
    cached_headlines = load_cached_data("news", "headlines") if cached_date == today_str else None

    if st.session_state.is_online:
        if not cached_headlines:
            headlines = get_news_from_rss()
            if headlines:
                for news in headlines:
                    st.markdown(f"**[{news['title']}]({news['link']})**")
                cache_data(headlines, "news", "headlines")
                st.session_state['news_cache_date'] = today_str
            else:
                st.error("❌ No headlines available")
        else:
            st.info("📰 Cached Headlines (Today):")
            for news in cached_headlines:
                if isinstance(news, dict) and 'title' in news and 'link' in news:
                    st.markdown(f"**[{news['title']}]({news['link']})**")
                else:
                    st.write(news)
    else:
        cached_headlines = get_offline_news()
        if cached_headlines:
            st.info("📰 Cached Headlines:")
            for news in cached_headlines:
                if isinstance(news, dict) and 'title' in news and 'link' in news:
                    st.markdown(f"**[{news['title']}]({news['link']})**")
                else:
                    st.write(news)
        else:
            st.warning("❌ No cached news data available")

# Code Sandbox Tab
with tabs[7]:
    st.subheader("🧪 Live Code Sandbox")
    st.info("Run Python code safely with timeout protection and restricted imports.")
    col1, col2 = st.columns([3, 1])
    with col1:
        code_input = st.text_area("Enter Python code:", height=300)
    with col2:
        st.markdown("**Available Modules:**")
        st.code("""
• math
• random  
• datetime
• json
• re
• collections
• itertools
• functools
• operator
• string
• textwrap
        """)
        timeout_setting = st.slider("Timeout (seconds)", 1, 30, 10)
    if st.button("🚀 Run Code", type="primary") and code_input.strip():
        with st.spinner("Executing code..."):
            result = execute_python_code(code_input, timeout=timeout_setting)
            if result['success']:
                st.success("✅ Code executed successfully!")
                if result['output']:
                    st.subheader("📤 Output:")
                    st.code(result['output'], language="text")
                if result['errors']:
                    st.warning("⚠️ Warnings/Errors:")
                    st.code(result['errors'], language="text")
            else:
                st.error("❌ Code execution failed!")
                st.code(result['errors'], language="text")

# Offline Mode Tab
with tabs[8]:
    st.subheader("Offline Mode Management")
    col1, col2, col3 = st.columns(3)
    with col1:
        if st.button("🔄 Check Connection"):
            st.session_state.is_online = check_internet_connection()
            st.rerun()
    with col2:
        status_color = "🟢" if st.session_state.is_online else "🔴"
        status_text = "Online" if st.session_state.is_online else "Offline"
        st.metric("Connection Status", f"{status_color} {status_text}")
    with col3:
        cache_dir = os.path.join(SAVE_DIR, "cache")
        cache_size = 0
        if os.path.exists(cache_dir):
            for root, dirs, files in os.walk(cache_dir):
                cache_size += sum(os.path.getsize(os.path.join(root, f)) for f in files)
        st.metric("Cache Size", f"{cache_size / 1024:.1f} KB")
    st.markdown("---")
    st.subheader("💾 Cached Data")
    cache_types = ["weather", "news", "documents", "images"]
    for cache_type in cache_types:
        cache_path = os.path.join(SAVE_DIR, "cache", cache_type)
        if os.path.exists(cache_path):
            files = os.listdir(cache_path)
            st.write(f"**{cache_type.title()}:** {len(files)} items cached")
            if files and st.checkbox(f"Show {cache_type} details"):
                for file in files[:5]:
                    file_path = os.path.join(cache_path, file)
                    file_size = os.path.getsize(file_path)
                    mod_time = datetime.fromtimestamp(os.path.getmtime(file_path))
                    st.write(f"  • {file} ({file_size} bytes, {mod_time.strftime('%Y-%m-%d %H:%M')})")
                if len(files) > 5:
                    st.write(f"  ... and {len(files) - 5} more")
        else:
            st.write(f"**{cache_type.title()}:** No cache available")

# Voice Features Tab
with tabs[9]:
    st.title("🎙️ Voice Features")
    st.header("🎧 Detect Emotion from Live Voice")
    if st.button("🎤 Start Recording (15 sec)"):
        audio_file = record_audio(duration=15)
        emotion = detect_emotion_simple(audio_file)
        st.audio(audio_file, format="audio/wav")
        st.success(f"🧠 Detected Emotion: **{emotion}**")
    st.markdown("---")
    st.header("🗨️ Type & Translate with Voice Reply")
    text_input = st.text_input("Enter your message here...")
    lang_map = {
        "English": "en", "Hindi": "hi", "Telugu": "te", "Tamil": "ta",
        "Kannada": "kn", "Malayalam": "ml", "Gujarati": "gu", "Marathi": "mr"
    }
    lang_selected = st.selectbox("Reply in language", list(lang_map.keys()))
    if text_input and lang_selected:
        lang_code = lang_map[lang_selected]
        translated_text = translate_text(text_input, lang_code)
        st.success(f"Translated ({lang_selected}): {translated_text}")
        speak_translated_text(translated_text, lang=lang_code)


# =========================
# 🛠️ AI Tools Tab
# =========================
with tabs[10]:
    st.title("🛠️ AI Tools")

   
    # 2. AI-powered Web Automation
    st.subheader("🤖 AI-powered Web Automation")
    st.write("Describe a web automation task (e.g., 'Download all images from https://example.com'). The bot will generate a Python script for you.")
    automation_prompt = st.text_area("Describe your automation task:")
    if st.button("Generate Automation Script"):
        if automation_prompt.strip():
            messages = [{"role": "user", "content": f"Write a Python script to: {automation_prompt}. Use only standard libraries and requests/bs4 if needed."}]
            response, model_used = get_ai_response(messages)
            st.markdown(f"**🤖 Response from {model_used}:**\n\n```python\n{response}\n```")
            st.code(response, language="python")

    st.markdown("---")

    # 3. AI-powered Resume Builder
    st.subheader("📄 AI-powered Resume Builder")
    st.write("Enter your full details below to generate a clean PDF resume with optional themes.")

    theme_color = st.selectbox("🎨 Select Theme", ["Default", "Light Blue", "Gray", "Dark"])
    theme_colors = {
        "Default": (0, 0, 0),
        "Light Blue": (70, 130, 180),
        "Gray": (80, 80, 80),
        "Dark": (30, 30, 30)
    }

    def sanitize_text(text):
        return text.encode('latin-1', 'replace').decode('latin-1')

    def build_resume_pdf(details, filename="generated_resume.pdf"):
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=10)
        pdf.set_font("Arial", 'B', 16)
        r, g, b = theme_colors[theme_color]
        pdf.set_text_color(r, g, b)
        pdf.cell(0, 10, sanitize_text(details["name"]), ln=True)
        pdf.set_font("Arial", '', 12)
        pdf.set_text_color(0, 0, 0)
        pdf.cell(0, 10, f"{sanitize_text(details['email'])} | {sanitize_text(details['phone'])}", ln=True)
        pdf.ln(5)

        def section(title, content_lines):
            pdf.set_font("Arial", 'B', 14)
            pdf.set_text_color(r, g, b)
            pdf.cell(0, 10, sanitize_text(title), ln=True)
            pdf.set_font("Arial", '', 12)
            pdf.set_text_color(0, 0, 0)
            for line in content_lines:
                pdf.multi_cell(0, 8, sanitize_text(f"- {line}"))
            pdf.ln(4)

        section("Profile Summary", [details["summary"]])
        section("Skills", [s.strip() for s in details["skills"].split(",") if s.strip()])
        section("Awards", [a.strip() for a in details["awards"].split("\n") if a.strip()])
        section("Organizations", [o.strip() for o in details["orgs"].split("\n") if o.strip()])
        section("Experience", [e.strip() for e in details["experience"].split("\n") if e.strip()])
        section("Job Profile", [details["job"]])
        section("Links", [l.strip() for l in details["links"].split("\n") if l.strip()])
        section("Education", [e.strip() for e in details["education"].split("\n") if e.strip()])

        pdf.output(filename)
        return filename

    # Input Fields
    name = st.text_input("Full Name")
    email = st.text_input("Email")
    phone = st.text_input("Phone Number")
    summary = st.text_area("Professional Summary")
    skills = st.text_area("Skills (comma separated)")
    awards = st.text_area("Awards & Achievements (one per line)")
    orgs = st.text_area("Organizations / Memberships (one per line)")
    links = st.text_area("Links (Portfolio, GitHub, LinkedIn, etc.) (one per line)")
    experience = st.text_area("Work Experience (one per line)")
    job = st.text_area("Job Profile / Objective")
    education = st.text_area("Education (one per line)")

    if st.button("📄 Generate Resume"):
        if name and email and phone:
            details = {
                "name": name,
                "email": email,
                "phone": phone,
                "summary": summary,
                "skills": skills,
                "awards": awards,
                "orgs": orgs,
                "links": links,
                "experience": experience,
                "job": job,
                "education": education
            }
            filename = build_resume_pdf(details, filename=f"{name.replace(' ', '_')}_resume.pdf")
            with open(filename, "rb") as f:
                st.download_button("📥 Download Your Resume", data=f, file_name=filename)
        else:
            st.error("⚠️ Please fill in at least your Name, Email, and Phone.")

    # 4. AI-powered Language Learning
    st.subheader("🌏 AI-powered Language Learning " + info_icon("Practice conversation, grammar correction, and vocabulary quizzes with AI."))
    st.write("Practice conversation, get grammar corrections, and vocabulary quizzes.")

    # Conversation Practice
    with st.container():
        st.markdown('<div class="glass-card" style="padding:16px;">', unsafe_allow_html=True)
        st.markdown("**Practice Conversation** " + info_icon("Type a message in your target language. The AI will correct and reply in that language."), unsafe_allow_html=True)
        lang_practice = st.text_area("Say something in English or your target language:", key="lang_practice")
        lang_target = st.selectbox("Target Language", ["English", "Hindi", "Telugu", "Tamil", "French", "Spanish", "German", "Chinese"], key="lang_target")
        if st.button("Practice Conversation", key="btn_practice_conv"):
            if lang_practice.strip():
                messages = [{"role": "user", "content": f"Correct my sentence and reply in {lang_target}: {lang_practice}"}]
                response, model_used = get_ai_response(messages)
                st.markdown(f"**🤖 Response from {model_used}:**\n\n{response}")
                speak_translated_text(response, lang="en" if lang_target=="English" else "hi" if lang_target=="Hindi" else "te" if lang_target=="Telugu" else "ta" if lang_target=="Tamil" else "en")
        st.markdown('</div>', unsafe_allow_html=True)

    # Explicit Grammar Correction
    with st.container():
        st.markdown('<div class="glass-card" style="padding:16px;">', unsafe_allow_html=True)
        st.markdown("**Grammar Correction** " + info_icon("Paste a sentence or paragraph. The AI will correct grammar, punctuation, and suggest improvements."), unsafe_allow_html=True)
        grammar_input = st.text_area("Enter text for grammar correction:", key="grammar_input")
        if st.button("Correct Grammar", key="btn_grammar_corr"):
            if grammar_input.strip():
                messages = [{"role": "user", "content": f"Correct all grammar, punctuation, and suggest improvements for: {grammar_input}"}]
                response, model_used = get_ai_response(messages)
                st.markdown(f"**📝 Corrected by {model_used}:**\n\n{response}")
        st.markdown('</div>', unsafe_allow_html=True)

    # Vocabulary Quiz (Enhanced)
    with st.container():
        st.markdown('<div class="glass-card" style="padding:16px;">', unsafe_allow_html=True)
        st.markdown("**Vocabulary Quiz** " + info_icon("Test your vocabulary with MCQ or fill-in-the-blank quizzes. Choose format and language."), unsafe_allow_html=True)
        quiz_word = st.text_input("Word for Vocabulary Quiz", key="quiz_word")
        quiz_format = st.radio("Quiz Format", ["Multiple Choice", "Fill in the Blank"], horizontal=True, key="quiz_format")
        if st.button("Get Quiz", key="btn_get_quiz"):
            if quiz_word.strip():
                prompt = f"Give me a {quiz_format.lower()} vocabulary quiz for the word '{quiz_word}' in {lang_target}. Provide the answer and a short explanation."
                messages = [{"role": "user", "content": prompt}]
                response, model_used = get_ai_response(messages)
                st.markdown(f"**🤖 Quiz from {model_used}:**\n\n{response}")
        st.markdown('</div>', unsafe_allow_html=True)

    st.markdown("---")

    # 5. AI-powered Social Media Post Generator
    st.subheader("📱 AI-powered Social Media Post Generator")
    platform = st.selectbox("Platform", ["Twitter", "LinkedIn", "Instagram", "Facebook"])
    post_topic = st.text_input("Post Topic or Idea")
    if st.button("Generate Post"):
        if post_topic.strip():
            messages = [{"role": "user", "content": f"Write a viral {platform} post about: {post_topic}. Add hashtags and emojis."}]
            response, model_used = get_ai_response(messages)
            st.markdown(f"**🤖 {platform} Post from {model_used}:**\n\n{response}")

    st.markdown("---")

    # 6. AI-powered Song Recommendation by Mood
    st.subheader("🎵 AI-powered Song Recommendation by Mood")

    mood = st.selectbox("Your Mood", ["Happy", "Sad", "Energetic", "Relaxed", "Romantic", "Motivated", "Party", "Chill"])
    search_term = st.text_input("🔍 Optional: Search songs or keyword (e.g., 'RRR', '2024 hits', 'love')")
    add_effect = st.checkbox("🎧 Add Reverb / Slowed Effect")
    num_results = st.slider("How many songs?", 5, 15, 10)

    if st.button("🎶 Recommend Songs"):
        st.info(f"🔎 Searching for latest + old Telugu, Hindi, English songs for mood: **{mood}** and keyword: **{search_term or 'N/A'}**...")

        try:
            from duckduckgo_search import DDGS
            from pydub import AudioSegment
            import tempfile

            playable_urls = []
            seen_titles = set()
            with DDGS() as ddgs:
                queries = [
                    f"{lang} {mood} {search_term} songs mp3 download site:archive.org"
                    for lang in ["Telugu", "Hindi", "English"]
                ]
                for q in queries:
                    for r in ddgs.text(q, max_results=10):
                        title = r.get("title", "").strip()
                        url = r.get("href", "")
                        if not title or title in seen_titles:
                            continue
                        if "archive.org" in url and (".mp3" in url or "/details/" in url):
                            playable_urls.append((title, url))
                            seen_titles.add(title)

            if not playable_urls:
                # Fallback: Curated local recommendations
                curated_songs = {
                    "Happy": [
                        {"title": "Happy", "artist": "Pharrell Williams", "img": "https://upload.wikimedia.org/wikipedia/en/9/92/Pharrell_Williams_-_Happy.jpg", "url": "https://www.soundhelix.com/examples/mp3/SoundHelix-Song-1.mp3"},
                        {"title": "Best Day Of My Life", "artist": "American Authors", "img": "https://upload.wikimedia.org/wikipedia/en/6/6b/American_Authors_Best_Day_of_My_Life.jpg", "url": "https://www.soundhelix.com/examples/mp3/SoundHelix-Song-2.mp3"}
                    ],
                    "Sad": [
                        {"title": "Someone Like You", "artist": "Adele", "img": "https://upload.wikimedia.org/wikipedia/en/9/9b/Adele_-_Someone_Like_You.png", "url": "https://www.soundhelix.com/examples/mp3/SoundHelix-Song-3.mp3"}
                    ],
                    "Energetic": [
                        {"title": "Stronger", "artist": "Kanye West", "img": "https://upload.wikimedia.org/wikipedia/en/6/6c/Kanye_West_-_Stronger.jpg", "url": "https://www.soundhelix.com/examples/mp3/SoundHelix-Song-4.mp3"}
                    ],
                    "Relaxed": [
                        {"title": "Weightless", "artist": "Marconi Union", "img": "https://upload.wikimedia.org/wikipedia/en/3/3c/Marconi_Union_Weightless.jpg", "url": "https://www.soundhelix.com/examples/mp3/SoundHelix-Song-5.mp3"}
                    ],
                    "Romantic": [
                        {"title": "Perfect", "artist": "Ed Sheeran", "img": "https://upload.wikimedia.org/wikipedia/en/4/45/Ed_Sheeran_Perfect_Single_cover.jpg", "url": "https://www.soundhelix.com/examples/mp3/SoundHelix-Song-6.mp3"}
                    ],
                    "Motivated": [
                        {"title": "Eye of the Tiger", "artist": "Survivor", "img": "https://upload.wikimedia.org/wikipedia/en/9/9e/Survivor_Eye_of_the_Tiger_single_cover.jpg", "url": "https://www.soundhelix.com/examples/mp3/SoundHelix-Song-7.mp3"}
                    ],
                    "Party": [
                        {"title": "Uptown Funk", "artist": "Mark Ronson ft. Bruno Mars", "img": "https://upload.wikimedia.org/wikipedia/en/b/b7/Mark_Ronson_-_Uptown_Funk_%28feat._Bruno_Mars%29.png", "url": "https://www.soundhelix.com/examples/mp3/SoundHelix-Song-8.mp3"}
                    ],
                    "Chill": [
                        {"title": "Sunset Lover", "artist": "Petit Biscuit", "img": "https://upload.wikimedia.org/wikipedia/en/7/7d/Petit_Biscuit_-_Sunset_Lover.png", "url": "https://www.soundhelix.com/examples/mp3/SoundHelix-Song-9.mp3"}
                    ]
                }
                mood_songs = curated_songs.get(mood, [])
                if mood_songs:
                    st.info("🎶 No online results. Showing curated recommendations:")
                    for song in mood_songs[:num_results]:
                        st.markdown(f"""**{song['title']}**  
Artist: {song['artist']}""")
                        st.audio(song["url"], format="audio/mp3")
                else:
                    st.warning("No curated songs available for this mood.")
            for title, url in playable_urls[:num_results]:
                try:
                    mp3_data = requests.get(url).content
                    temp_path = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3").name
                    with open(temp_path, "wb") as f:
                        f.write(mp3_data)
                    if add_effect:
                        audio = AudioSegment.from_file(temp_path)
                        slowed = audio.speedup(playback_speed=0.8).fade_in(500).fade_out(500)
                        effect_path = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3").name
                        slowed.export(effect_path, format="mp3")
                        st.markdown(f"🎧 **{title}** (Slowed)")
                        st.audio(effect_path, format="audio/mp3")
                    else:
                        st.markdown(f"🎵 **{title}**")
                        st.audio(temp_path, format="audio/mp3")
                    shown += 1
                except:
                    continue
        except Exception as e:
            st.error(f"❌ Failed to fetch songs: {e}")


# =============================
# 🎮 Interactive Game Center Tab
# =============================
# =============================
# 🎮 Ultimate Game Center Tab
# =============================
with tabs[11]:  # Adjust index if needed
    st.title("🎮 Interactive Game Center")

    game = st.selectbox("Choose a game to play:", [
        "🧠 Riddle of the Day",
        "❓ Trivia Quiz",
        "🔤 Word Scramble",
        "⭕ Tic-Tac-Toe (2-Player)",
        "✊ Rock-Paper-Scissors",
        "🧠 Emoji Memory Match",
        "🧮 2048 Lite",
        "🐹 Whack-a-Mole",
        "🌍 Guess the Flag",
        "⚡ Reaction Speed Test"
    ])

    import time
    import random

    # Game 1: Riddle of the Day
    if game == "🧠 Riddle of the Day":
        riddles = {
            "What has keys but can't open locks?": "keyboard",
            "The more you take, the more you leave behind. What am I?": "footsteps",
            "What comes once in a minute, twice in a moment, but never in a thousand years?": "m"
        }
        question = list(riddles.keys())[datetime.now().day % len(riddles)]
        st.markdown(f"**🧩 Riddle:** {question}")
        answer = st.text_input("Your answer:")
        if answer:
            if answer.lower().strip() == riddles[question]:
                st.success("✅ Correct!")
            else:
                st.error("❌ Try again!")

    # Game 2: Trivia Quiz
    elif game == "❓ Trivia Quiz":
        trivia_q = {
            "Which planet is known as the Red Planet?": ("Mars", ["Earth", "Venus", "Mars", "Jupiter"]),
            "Who wrote Hamlet?": ("Shakespeare", ["Shakespeare", "Charles Dickens", "Tolstoy", "Homer"]),
            "What is the capital of Japan?": ("Tokyo", ["Seoul", "Beijing", "Tokyo", "Bangkok"]),
        }
        q = list(trivia_q.keys())[datetime.now().day % len(trivia_q)]
        correct, options = trivia_q[q]
        st.markdown(f"**🤔 Question:** {q}")
        choice = st.radio("Choose an answer:", options)
        if st.button("Submit Answer"):
            if choice == correct:
                st.success("🎉 That's correct!")
            else:
                st.error(f"Oops! Correct answer was: **{correct}**")

    # Game 3: Word Scramble
    elif game == "🔤 Word Scramble":
        words = ["streamlit", "chatbot", "python", "response", "prompt"]
        word = random.choice(words)
        scrambled = ''.join(random.sample(word, len(word)))
        st.markdown(f"**🔁 Unscramble this word:** `{scrambled}`")
        user_guess = st.text_input("Your guess:")
        if user_guess:
            if user_guess.lower().strip() == word:
                st.success("🎯 Well done!")
            else:
                st.warning("Not quite, try again!")

    # Game 4: Tic-Tac-Toe
    elif game == "⭕ Tic-Tac-Toe (2-Player)":
        if "ttt_board" not in st.session_state:
            st.session_state.ttt_board = [""] * 9
            st.session_state.ttt_turn = "X"

        def check_winner(b):
            wins = [(0,1,2),(3,4,5),(6,7,8),(0,3,6),(1,4,7),(2,5,8),(0,4,8),(2,4,6)]
            for i,j,k in wins:
                if b[i] and b[i] == b[j] == b[k]:
                    return b[i]
            return None

        cols = st.columns(3)
        for i in range(9):
            if cols[i % 3].button(st.session_state.ttt_board[i] or " ", key=f"ttt_{i}"):
                if not st.session_state.ttt_board[i]:
                    st.session_state.ttt_board[i] = st.session_state.ttt_turn
                    winner = check_winner(st.session_state.ttt_board)
                    if winner:
                        st.success(f"🎉 Player {winner} wins!")
                    elif all(st.session_state.ttt_board):
                        st.info("🤝 It's a draw!")
                    st.session_state.ttt_turn = "O" if st.session_state.ttt_turn == "X" else "X"

        if st.button("🔁 Reset Game"):
            st.session_state.ttt_board = [""] * 9
            st.session_state.ttt_turn = "X"

    # Game 5: Rock-Paper-Scissors
    elif game == "✊ Rock-Paper-Scissors":
        options = ["Rock", "Paper", "Scissors"]
        user_choice = st.selectbox("Your move:", options)
        if st.button("Play"):
            bot = random.choice(options)
            st.write(f"🤖 Bot chose: **{bot}**")
            if user_choice == bot:
                st.info("😐 It's a tie!")
            elif (user_choice == "Rock" and bot == "Scissors") or \
                 (user_choice == "Paper" and bot == "Rock") or \
                 (user_choice == "Scissors" and bot == "Paper"):
                st.success("🎉 You win!")
            else:
                st.error("😢 You lose.")

    # Game 6: Emoji Memory Match
    elif game == "🧠 Emoji Memory Match":
        emojis = ["🐶", "🐱", "🐶", "🐱"]
        if "mem_board" not in st.session_state:
            random.shuffle(emojis)
            st.session_state.mem_board = emojis
            st.session_state.mem_flips = [False]*4
            st.session_state.mem_prev = -1

        cols = st.columns(2)
        for i in range(4):
            with cols[i % 2]:
                label = st.session_state.mem_board[i] if st.session_state.mem_flips[i] else "❓"
                if st.button(label, key=f"mem_{i}"):
                    if not st.session_state.mem_flips[i]:
                        st.session_state.mem_flips[i] = True
                        if st.session_state.mem_prev == -1:
                            st.session_state.mem_prev = i
                        else:
                            j = st.session_state.mem_prev
                            if st.session_state.mem_board[i] != st.session_state.mem_board[j]:
                                st.warning("❌ Not a match!")
                                st.session_state.mem_flips[i] = False
                                st.session_state.mem_flips[j] = False
                            st.session_state.mem_prev = -1
        if all(st.session_state.mem_flips):
            st.success("🎉 You matched all!")
        if st.button("🔁 Restart"):
            random.shuffle(emojis)
            st.session_state.mem_board = emojis
            st.session_state.mem_flips = [False]*4
            st.session_state.mem_prev = -1

    # Game 7: 2048 Lite
    elif game == "🧮 2048 Lite":
        st.info("🚧 2048 Game Placeholder: Full grid version coming soon!")
        st.write("Hint: You can later integrate keyboard events + NumPy to make a full clone.")

    # Game 8: Whack-a-Mole (Text)
    elif game == "🐹 Whack-a-Mole":
        if st.button("Start Game"):
            mole = random.randint(1, 9)
            st.session_state.mole = mole
        mole = st.session_state.get("mole", 0)
        cols = st.columns(3)
        for i in range(9):
            if cols[i % 3].button("🐹" if i+1 == mole else "🥱", key=f"mole_{i}"):
                if i+1 == mole:
                    st.success("Bonk! 🎯")
                else:
                    st.warning("Missed!")

    # Game 9: Guess the Flag
    elif game == "🌍 Guess the Flag":
        flags = {
            "🇯🇵": "Japan",
            "🇫🇷": "France",
            "🇧🇷": "Brazil",
            "🇮🇳": "India",
            "🇺🇸": "United States"
        }
        emoji, country = random.choice(list(flags.items()))
        st.markdown(f"**🏁 Flag:** {emoji}")
        guess = st.text_input("Which country is this?")
        if guess:
            if guess.lower() == country.lower():
                st.success("🎉 Correct!")
            else:
                st.error(f"❌ It's actually {country}.")

    # Game 10: Reaction Speed Test
    elif game == "⚡ Reaction Speed Test":
        if "react_started" not in st.session_state:
            st.session_state.react_started = False
            st.session_state.react_start_time = 0

        if not st.session_state.react_started:
            if st.button("Start Reaction Test"):
                wait_time = random.randint(2, 5)
                st.session_state.react_started = True
                time.sleep(wait_time)
                st.session_state.react_start_time = time.time()
                st.session_state.react_ready = True
                st.rerun()
        else:
            if st.button("CLICK NOW!"):
                reaction = time.time() - st.session_state.react_start_time
                st.success(f"⚡ Your reaction time: {reaction:.3f} seconds")
                st.session_state.react_started = False



# =============================
# 🩺 Health Assistant Tab
# =============================
with tabs[12]:
    st.title("🩺 Health Bot - Symptom Checker")

    symptoms = st.multiselect("Select symptoms you're experiencing:", [
        "Fever", "Headache", "Cold", "Cough", "Sore throat", "Body pain", 
        "Eye strain", "Stomach ache", "Nausea", "Fatigue"
    ])

    if st.button("🧠 Get Remedies"):
        if not symptoms:
            st.warning("Please select at least one symptom.")
        else:
            st.success("✅ Based on your symptoms, here are suggestions:")

            remedies = []
            tablets = []

            if "Fever" in symptoms:
                remedies.append("Drink tulsi + ginger tea 2 times daily.")
                tablets.append("Paracetamol 500mg (if fever >100.5°F)")

            if "Cold" in symptoms or "Sore throat" in symptoms:
                remedies.append("Gargle warm salt water. Inhale steam with eucalyptus oil.")
                tablets.append("Cetirizine (for sneezing or runny nose)")

            if "Headache" in symptoms or "Eye strain" in symptoms:
                remedies.append("Rest eyes. Apply cold pack. Drink water.")
                tablets.append("Dolo 650 (mild pain relief)")

            if "Stomach ache" in symptoms or "Nausea" in symptoms:
                remedies.append("Jeera water. Avoid spicy food. Eat light khichdi.")
                tablets.append("Domperidone (nausea), Meftal Spas (cramps)")

            if "Fatigue" in symptoms or "Body pain" in symptoms:
                remedies.append("Rest + hydrate. Banana or dates for quick energy.")
                tablets.append("DOLO 650 or Ibuprofen")

            st.markdown("### 🌿 Ayurvedic / Home Remedies")
            for r in remedies:
                st.markdown(f"- {r}")

            st.markdown("### 💊 Optional Tablets")
            for t in tablets:
                st.markdown(f"- {t} *(Only if needed. Consult a doctor if unsure.)*")


# Export, Logout, Delete Account section
st.markdown("---")
st.header("Export & Account Management")
col1, col2, col3 = st.columns(3)
with col1:
    if st.button("📤 Export My Chats"):
        c.execute("SELECT * FROM chat_history WHERE email=?", (st.session_state.email,))
        data = c.fetchall()
        filename = f"chat_history_{st.session_state.email}.txt"
        with open(filename, "w", encoding="utf-8") as f:
            for row in data:
                f.write(f"{row[4]} - You: {row[2]} | Bot: {row[3]}\n")
        with open(filename, "rb") as f:
            st.download_button("Download Chat History", data=f, file_name=filename)
with col2:
    if st.button("🚪 Logout"):
        del st.session_state.email
        st.rerun()
with col3:
    if st.button("🗑️ Delete My Account"):
        c.execute("DELETE FROM users WHERE email=?", (st.session_state.email,))
        c.execute("DELETE FROM chat_history WHERE email=?", (st.session_state.email,))
        conn.commit()
        st.success("Account deleted successfully.")
        del st.session_state.email
        st.rerun()
